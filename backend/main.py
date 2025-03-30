import os
import re
from dotenv import load_dotenv
import requests
from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.responses import FileResponse
from pydantic import BaseModel
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document
import pandas as pd
import google.generativeai as genai

# Load API key securely
load_dotenv()
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
genai.configure(api_key=GEMINI_API_KEY)

# FastAPI app setup
app = FastAPI()

UPLOAD_DIR = "uploaded_files"
os.makedirs(UPLOAD_DIR, exist_ok=True)

# Input data model for direct text input
class RequirementInput(BaseModel):
    text: str

# Home Route
@app.get("/")
async def home():
    return {"message": "Welcome to the Requirement Extraction API!"}

# File Upload and Extraction Route (Combined improved version)
@app.post("/upload/")
async def upload_file(file: UploadFile = File(...)):
    file_path = os.path.join(UPLOAD_DIR, file.filename)
    print(f"📂 Received file: {file.filename}")

    try:
        # Save the uploaded file
        with open(file_path, "wb") as buffer:
            buffer.write(await file.read())
        print(f"✅ File saved at: {file_path}")
    except Exception as e:
        print(f"❌ Error saving file: {str(e)}")
        raise HTTPException(status_code=500, detail="File processing failed")

    # Extract text
    extracted_text = extract_text(file_path)
    print(f"📝 Extracted text:\n{extracted_text[:200]}...")  # Show first 200 chars

    if extracted_text == "Unsupported file format":
        print("🚫 Unsupported file format")
        raise HTTPException(status_code=400, detail="Unsupported file format")

    # Get requirements from Gemini
    try:
        extracted_requirements = get_requirements_gemini(extracted_text)
        print(f"📌 Extracted Requirements:\n{extracted_requirements[:200]}...")
    except Exception as e:
        print(f"❌ Error extracting requirements: {str(e)}")
        raise HTTPException(status_code=500, detail="Error extracting requirements")

    # Generate documents
    try:
        word_file = generate_word_file(extracted_requirements)
        excel_file = generate_excel_file(extracted_requirements)
        
        return {
            "message": "Files generated successfully!",
            "word_file": word_file,
            "excel_file": excel_file,
            "requirements": extracted_requirements,
            "filename": file.filename
        }
    except Exception as e:
        print(f"❌ Error generating documents: {str(e)}")
        raise HTTPException(status_code=500, detail="Error generating documents")

# Text Extraction Logic (Combined from both)
def extract_text(file_path):
    if not os.path.exists(file_path):
        return f"Error: File '{file_path}' not found"

    file_ext = file_path.split(".")[-1].lower()

    try:
        if file_ext == "pdf":
            doc = fitz.open(file_path)
            return "\n".join([page.get_text("text") for page in doc])
        elif file_ext in ["doc", "docx"]:
            doc = Document(file_path)
            return "\n".join([p.text for p in doc.paragraphs])
        elif file_ext in ["xls", "xlsx"]:
            df = pd.read_excel(file_path)
            return df.to_string()
        elif file_ext == "pptx":
            prs = Presentation(file_path)
            return "\n".join([ 
                shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")
            ])
        elif file_ext == "txt":
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()
        else:
            return "Unsupported file format"
    except Exception as e:
        return f"Error extracting text: {str(e)}"

# Gemini API Integration (Combined approaches)
def get_requirements_gemini(input_text):
    prompt = f"""Extract and categorize the following text into Functional and Non-Functional requirements.
    Format your response with clear headings:
    **Functional Requirements**
    - [List functional requirements as bullet points]
    
    **Non-Functional Requirements**
    - [List non-functional requirements as bullet points]
    
    Here is the input text:
    {input_text}"""

    try:
        # Try both approaches (direct API call and genai library)
        try:
            model = genai.GenerativeModel("gemini-2.0-flash-exp")
            response = model.generate_content(prompt)
            return response.text
        except:
            url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash:generateContent?key={GEMINI_API_KEY}"
            headers = {'Content-Type': 'application/json'}
            data = {"contents": [{"parts": [{"text": prompt}]}]}
            response = requests.post(url, json=data, headers=headers)
            if response.status_code == 200:
                return response.json()['candidates'][0]['content']['parts'][0]['text']
            raise Exception("Gemini API error")
    except Exception as e:
        print(f"Gemini API Error: {str(e)}")
        return f"Error processing requirements: {str(e)}"

# Document Generation Functions (Improved)
def generate_word_file(extracted_requirements):
    # Parse requirements into sections
    functional, non_functional = parse_requirements(extracted_requirements)
    
    doc = Document()
    doc.add_heading('Extracted Requirements', 0)
    
    # Functional Requirements
    doc.add_heading('Functional Requirements', level=1)
    for req in functional.split('\n'):
        if req.strip():
            doc.add_paragraph(req.strip(), style="List Bullet")
    
    # Non-Functional Requirements
    doc.add_heading('Non-Functional Requirements', level=1)
    for req in non_functional.split('\n'):
        if req.strip():
            doc.add_paragraph(req.strip(), style="List Bullet")

    # Save the document
    word_file = "requirements.docx"
    doc.save(word_file)
    return word_file

def parse_requirements(extracted_text):
    # Improved parsing with regex
    functional_match = re.search(r"\*\*Functional Requirements\*\*(.*?)(\*\*Non-Functional Requirements\*\*|$)", 
                                extracted_text, re.DOTALL | re.IGNORECASE)
    non_functional_match = re.search(r"\*\*Non-Functional Requirements\*\*(.*)", 
                                   extracted_text, re.DOTALL | re.IGNORECASE)
    
    functional = functional_match.group(1).strip() if functional_match else "No functional requirements found."
    non_functional = non_functional_match.group(1).strip() if non_functional_match else "No non-functional requirements found."
    
    return functional, non_functional

def generate_excel_file(requirements):
    functional, non_functional = parse_requirements(requirements)
    
    # Create a DataFrame with requirements
    data = {
        "Type": ["Functional"] * len(functional.split('\n')) + ["Non-Functional"] * len(non_functional.split('\n')),
        "Requirement": [req.strip() for req in functional.split('\n') if req.strip()] + 
                      [req.strip() for req in non_functional.split('\n') if req.strip()],
        "Priority": ["TBD"] * (len(functional.split('\n')) + len(non_functional.split('\n')))
    }
    
    df = pd.DataFrame(data)
    file_path = "user_stories.xlsx"
    df.to_excel(file_path, index=False)
    return file_path

# File Download Routes
@app.get("/download-word/")
async def download_word():
    return FileResponse("requirements.docx", 
                      media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document", 
                      filename="requirements.docx")

@app.get("/download-excel/")
async def download_excel():
    return FileResponse("user_stories.xlsx", 
                      media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", 
                      filename="user_stories.xlsx")

# Requirement Extraction Endpoint (For Direct Text)
@app.post("/extract-requirements/")
async def extract_requirements(input_data: RequirementInput):
    extracted_requirements = get_requirements_gemini(input_data.text)
    
    # Generate documents for direct text input too
    word_file = generate_word_file(extracted_requirements)
    excel_file = generate_excel_file(extracted_requirements)
    
    return {
        "requirements": extracted_requirements,
        "word_file": word_file,
        "excel_file": excel_file
    }

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)